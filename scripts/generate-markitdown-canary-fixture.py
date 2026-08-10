#!/usr/bin/env python3
"""Generate redaction-safe binary fixtures for the MarkItDown live canary."""

from __future__ import annotations

import argparse
import base64
import math
import struct
import wave
from pathlib import Path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-type", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--payload-base64", required=True)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = Path(args.output)
    payload = base64.b64decode(args.payload_base64).decode("utf-8")
    output.parent.mkdir(parents=True, exist_ok=True)

    generators = {
        "pdf": generate_pdf,
        "docx": generate_docx,
        "pptx": generate_pptx,
        "xlsx": generate_xlsx,
        "image": generate_image,
        "audio": generate_audio,
    }
    generator = generators.get(args.source_type)
    if generator is None:
        raise ValueError(f"unsupported binary fixture type: {args.source_type}")
    generator(output, payload)


def generate_pdf(output: Path, payload: str) -> None:
    from reportlab.pdfgen import canvas

    document = canvas.Canvas(str(output))
    y = 800
    for line in payload.splitlines():
        document.drawString(72, y, line[:110])
        y -= 16
        if y < 72:
            document.showPage()
            y = 800
    document.save()


def generate_docx(output: Path, payload: str) -> None:
    from docx import Document

    document = Document()
    for line in payload.splitlines():
        document.add_paragraph(line)
    document.save(output)


def generate_pptx(output: Path, payload: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    frame = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(8.8), Inches(6.2)).text_frame
    frame.word_wrap = True
    for index, line in enumerate(payload.splitlines()):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
    presentation.save(output)


def generate_xlsx(output: Path, payload: str) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    for row, line in enumerate(payload.splitlines(), start=1):
        sheet.cell(row, 1, line)
    workbook.save(output)


def generate_image(output: Path, payload: str) -> None:
    from PIL import Image, ImageDraw, PngImagePlugin

    metadata = PngImagePlugin.PngInfo()
    metadata.add_text("Description", payload)
    image = Image.new("RGB", (1200, 900), "white")
    ImageDraw.Draw(image).multiline_text((40, 40), payload, fill="black", spacing=8)
    image.save(output, pnginfo=metadata)


def generate_audio(output: Path, payload: str) -> None:
    # The payload is intentionally not encoded into samples. MarkItDown's
    # transcription boundary must prove speech semantics through an approved
    # provider rather than fixture metadata.
    del payload
    sample_rate = 16_000
    with wave.open(str(output), "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(sample_rate)
        frames = bytearray()
        for index in range(sample_rate):
            sample = int(6_000 * math.sin(2 * math.pi * 440 * index / sample_rate))
            frames.extend(struct.pack("<h", sample))
        audio.writeframes(frames)


if __name__ == "__main__":
    main()
